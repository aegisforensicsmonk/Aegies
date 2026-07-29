from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    
    # Customize title font
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 51, 102) # Dark blue
        
def add_bullet_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = title_text
    
    # Customize title
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(0, 51, 102)
    
    tf = body.text_frame
    tf.clear() # clear default paragraphs
    
    for point in bullet_points:
        p = tf.add_paragraph()
        if isinstance(point, tuple):
            # It's a sub-bullet (indent level 1)
            p.text = point[0]
            p.level = 1
            p.font.size = Pt(20)
        else:
            p.text = point
            p.level = 0
            p.font.bold = True if ":" in point else False
            p.font.size = Pt(24)

prs = Presentation()

# Slide 1
add_title_slide(prs, "CIIP: Cyber Intelligence & Investigation Platform", "Next-Generation Digital Forensics Hub for I4C\n\nPresented by: Aegis Forensics Monk")

# Slide 2
add_bullet_slide(prs, "The Challenge in Modern Cyber Investigations", [
    "Fragmented Investigative Tools:", 
    ("Investigators juggle disjointed software for OSINT, IPDR, and Malware analysis, leading to context switching and lost intelligence.",),
    "Chain of Custody Risks:",
    ("Maintaining strict, legally admissible chain of custody across massive digital evidence sets is highly prone to manual errors.",),
    "Reporting Bottlenecks:",
    ("Drafting comprehensive, court-ready forensic reports consumes hours that should be spent on active investigation.",),
    "Data Overload & Silos:",
    ("Correlating disparate datasets (IPs, Call Data Records, Cryptocurrency hashes) without automated visualization is highly inefficient.",)
])

# Slide 3
add_bullet_slide(prs, "The Solution: Introducing CIIP", [
    "Unified Intelligence Dashboard:",
    ("A single, cohesive platform combining OSINT, IPDR tracking, and Ransomware analysis into one intuitive interface.",),
    "Built for Law Enforcement:",
    ("Designed specifically around the operational needs and legal constraints of authorized investigators.",),
    "Security & Integrity First:",
    ("Engineered from the ground up to ensure absolute evidence admissibility and platform security.",),
    "Force Multiplier:",
    ("Acts as a centralized hub that accelerates investigations and drastically reduces case turnaround time.",)
])

# Slide 4
add_bullet_slide(prs, "Core Investigative Modules (1/2)", [
    "Advanced OSINT Lab:",
    ("Automated open-source intelligence gathering for IPs, Domains, Emails, and Hashes.",),
    ("Includes automated confidence scoring to prioritize high-value leads.",),
    "IPDR / CDR Analyzer:",
    ("In-depth analysis of IP Detail Records and Call Data Records.",),
    ("Detects communication patterns, frequencies, and maps suspect geolocations automatically.",)
])

# Slide 5
add_bullet_slide(prs, "Core Investigative Modules (2/2)", [
    "Ransomware Analysis Lab:",
    ("Dedicated secure environment for IOC (Indicators of Compromise) extraction and analysis.",),
    ("Active tracking of ransomware threats and related artifacts.",),
    "Interactive Case Management:",
    ("Real-time intelligence overview with active case tracking.",),
    ("Visual Timelines and Entity Relationship Graphs (via Cytoscape.js) to map criminal networks.",)
])

# Slide 6
add_bullet_slide(prs, "Unbreakable Chain of Custody", [
    "Immutable Evidence Vault:",
    ("Evidence cannot be altered or deleted (soft-delete only policies applied).",),
    "Automated Cryptographic Hashing:",
    ("SHA-256 hashes are automatically computed the moment evidence is uploaded.",),
    "Tamper-Proof Audit Trails:",
    ("Database-level triggers in PostgreSQL guarantee audit logs cannot be modified, bypassing application-level vulnerabilities.",),
    "Strict Role-Based Access Control (RBAC):",
    ("Granular permissions (Admin, Investigator, Analyst, Supervisor) secured via JWT.",)
])

# Slide 7
add_bullet_slide(prs, "AI-Assisted Intelligence & Reporting", [
    "Automated Report Drafting:",
    ("Utilizes advanced AI (OpenAI-compatible or secure Local LLM) to instantly generate template-based forensic reports.",),
    "Complex Data Summarization:",
    ("AI synthesizes timelines, evidence correlations, and massive datasets into digestible executive summaries.",),
    "Strict Advisory Guardrails:",
    ("All AI outputs are clearly flagged as 'Advisory'.",),
    ("The human investigator always remains the final authority, ensuring judicial standards are never compromised.",)
])

# Slide 8
add_bullet_slide(prs, "Robust & Scalable Architecture", [
    "Frontend (The Interface):",
    ("Next.js 14+ (React) with Tailwind CSS for a modern, responsive UI.",),
    ("Integrated mapping (Leaflet) and graphing tools.",),
    "Backend (The Brain):",
    ("High-performance Python FastAPI handling secure REST endpoints.",),
    ("Asynchronous background task processing via RabbitMQ.",),
    "Infrastructure (The Engine Room):",
    ("PostgreSQL 16: Immutable structured data storage.",),
    ("Neo4j: Graph database for suspect/IP relationship mapping.",),
    ("Elasticsearch 8 & MinIO: Lightning-fast search and secure S3-compatible evidence storage.",)
])

# Slide 9
add_bullet_slide(prs, "Deployment & Operational Readiness", [
    "Fully Containerized Environment:",
    ("The entire stack is containerized using Docker.",),
    "Rapid & Standardized Deployment:",
    ("Can be spun up in minutes using 'docker-compose up -d'.",),
    "Air-Gapped Network Capable:",
    ("Architecture is designed to function seamlessly within highly secure, isolated law enforcement networks.",),
    "Highly Scalable:",
    ("Built to scale from a local precinct deployment to a centralized national database.",)
])

# Slide 10
add_bullet_slide(prs, "Why CIIP is the Ideal Fit for I4C", [
    "Centralized Cyber Intelligence:",
    ("Perfectly aligns with I4C's mandate to coordinate complex cybercrime investigations across states.",),
    "Efficiency Multiplier:",
    ("Drastically reduces the timeline from evidence collection to actionable intelligence.",),
    "Guaranteed Court-Ready Evidence:",
    ("Ensures all digital artifacts meet the rigorous legal standards of the Indian judicial system.",),
    "Future-Proof Investigation Platform:",
    ("Seamlessly integrates cutting-edge AI while maintaining traditional forensic rigor.",)
])

# Slide 11
add_bullet_slide(prs, "Conclusion & Next Steps", [
    "Thank You for Your Time",
    ("Opening the floor for Questions.",),
    "Live Demonstration Available",
    ("We can now showcase a live walkthrough of the CIIP dashboard, OSINT Lab, and Evidence Vault.",)
])

prs.save('CIIP_I4C_Detailed_Presentation.pptx')
print("Successfully generated CIIP_I4C_Detailed_Presentation.pptx")
